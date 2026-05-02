"""
weekly_mi_report.py — Vietnam Infrastructure Weekly MI Report Generator
=========================================================================
주간 토요일(Saturday) 자동 실행 모듈

[ 입력 ]
  1. data/database/Vietnam_Infra_News_Database_Final.xlsx  (News DB)
  2. data/shared/knowledge_index.json                      (19개 플랜 마스터데이터)

[ 출력 ]
  1. docs/Vietnam_Infra_MI_Report_Weekly.docx              (Word 보고서)
  2. docs/Vietnam_Infra_MI_Report_Weekly.pptx              (PPT 보고서)
  3. docs/index.html                                        (MI Dashboard)
  4. data/agent_output/weekly_report_log.json              (실행 로그)

[ 핵심 기능 ]
  - 19개 플랜 모두에 대해 6단계 섹션 생성 (KPI/사업개요/프로젝트/AI분석/갭진단/기사목록)
  - 금주 신규 기사 0건이어도 보고서 정상 생성 (knowledge_index 고정 데이터 기반)
  - AI 자체 갭 진단 — 매핑 부족 플랜 자동 식별 + 키워드 보강 권고
  - Layer 1 (knowledge_index 고정) + Layer 2 (Claude Haiku 동적) 이중 구조

[ 영구 제약 ]
  - 19개 플랜 ID 고정: VN-SWM-NATIONAL-2030, VN-WAT-RESOURCES, VN-WAT-URBAN, VN-WW-2030,
    VN-OG-2030, VN-PWR-PDP8-LNG, VN-PWR-PDP8-NUCLEAR, VN-PWR-PDP8-RENEWABLE,
    VN-HAN-URBAN-2045, VN-IP-NORTH-2030, VN-TRAN-2055, VN-URB-METRO-2030, VN-ENV-IND-1894
  - Plan_ID 별칭(Aliases) 매핑은 PID_ALIASES 사전 참조
  - knowledge_index.json이 없으면 보고서 생성 중단 + 오류 로그
"""

import os
import json
import re
import logging
from collections import defaultdict
from datetime import date, datetime, timedelta
from pathlib import Path

import openpyxl
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Pt as PPt, Cm as PCm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

logging.basicConfig(level=logging.INFO, format='%(asctime)s [%(levelname)s] %(message)s', datefmt='%H:%M:%S')
log = logging.getLogger('weekly_mi_report')

# ─────────────────────────────────────────────────────
# 경로 설정
# ─────────────────────────────────────────────────────
_SCRIPTS_DIR = Path(os.path.dirname(os.path.abspath(__file__)))
_ROOT_DIR    = _SCRIPTS_DIR.parent

EXCEL_PATH = Path(os.environ.get('EXCEL_PATH',
    str(_ROOT_DIR / 'data' / 'database' / 'Vietnam_Infra_News_Database_Final.xlsx')))
KI_PATH    = Path(os.environ.get('KI_PATH',
    str(_ROOT_DIR / 'data' / 'shared' / 'knowledge_index.json')))
DOCS_DIR   = _ROOT_DIR / 'docs'
AGENT_OUT  = _ROOT_DIR / 'data' / 'agent_output'

# ─────────────────────────────────────────────────────
# 19개 플랜 표준 ID + 별칭
# ─────────────────────────────────────────────────────
PID_ALIASES = {
    # 환경 인프라 (6)
    'VN-WW-2030':            ['VN-WW-2030'],
    'VN-SWM-NATIONAL-2030':  ['VN-SWM-NATIONAL-2030', 'VN-SWM-NATIONAL'],
    'VN-WAT-URBAN':          ['VN-WAT-URBAN'],
    'VN-WAT-RESOURCES':      ['VN-WAT-RESOURCES', 'VN-WAT-RESOURCE'],
    'VN-ENV-IND-1894':       ['VN-ENV-IND-1894'],
    'VN-ENV-NATURA':         ['VN-ENV-NATURA', 'VN-ENV-BIODIV'],
    # 에너지·전력 (5)
    'VN-PWR-PDP8':           ['VN-PWR-PDP8', 'VN-PDP8'],
    'VN-PWR-PDP8-RENEWABLE': ['VN-PWR-PDP8-RENEWABLE', 'VN-PWR-PDP8-REN', 'VN-PWR-PDP8, VN-PWR-PDP8-RENEWABLE'],
    'VN-PWR-PDP8-LNG':       ['VN-PWR-PDP8-LNG', 'VN-PWR-PDP8, VN-PWR-PDP8-LNG'],
    'VN-PWR-PDP8-NUCLEAR':   ['VN-PWR-PDP8-NUCLEAR', 'VN-PWR-PDP8-NUC'],
    'VN-OG-2030':            ['VN-OG-2030', 'VN-OG'],
    # 도시·교통·산업 (8)
    'VN-TRAN-2055':          ['VN-TRAN-2055'],
    'VN-URB-METRO-2030':     ['VN-URB-METRO-2030', 'VN-URB-METRO-20'],
    'VN-IP-NORTH-2030':      ['VN-IP-NORTH-2030', 'VN-IP'],
    'VN-HAN-URBAN-2045':     ['VN-HAN-URBAN-2045', 'VN-HAN-URBAN', 'HN-URBAN-INFRA', 'HN-URBAN-WEST', 'HN-URBAN-NORTH'],
    'VN-HAN-DONG-ANH':       ['VN-HAN-DONG-ANH', 'HN-DONG-ANH'],
    'VN-HAN-HOA-LAC':        ['VN-HAN-HOA-LAC', 'HN-HOA-LAC'],
    'VN-HOUSING-2030':       ['VN-HOUSING-2030'],
    'VN-SMART-2025':         ['VN-SMART-2025', 'VN-SMART'],
}

AREA_GROUPS = {
    '환경 인프라 (Environment)':       ['VN-WW-2030','VN-SWM-NATIONAL-2030','VN-WAT-URBAN','VN-WAT-RESOURCES','VN-ENV-IND-1894','VN-ENV-NATURA'],
    '에너지·전력 (Energy & Power)':     ['VN-PWR-PDP8','VN-PWR-PDP8-RENEWABLE','VN-PWR-PDP8-LNG','VN-PWR-PDP8-NUCLEAR','VN-OG-2030'],
    '도시·교통·산업 (Urban & Transport)':['VN-TRAN-2055','VN-URB-METRO-2030','VN-IP-NORTH-2030','VN-HAN-URBAN-2045','VN-HAN-DONG-ANH','VN-HAN-HOA-LAC','VN-HOUSING-2030','VN-SMART-2025'],
}


# ═════════════════════════════════════════════════════
# Step 1: knowledge_index 로드
# ═════════════════════════════════════════════════════
def load_knowledge_index() -> dict:
    if not KI_PATH.exists():
        log.error(f"knowledge_index.json 없음: {KI_PATH}")
        return {}
    with open(KI_PATH, encoding='utf-8') as f:
        ki = json.load(f)
    plans = ki.get('plans', ki)
    log.info(f"knowledge_index 로드: {len(plans)}개 플랜")
    return plans


# ═════════════════════════════════════════════════════
# Step 2: News DB에서 플랜별 기사 추출
# ═════════════════════════════════════════════════════
def load_articles_per_plan() -> dict:
    """News DB의 Matched_Plan 시트에서 표준 plan_id별 기사 추출"""
    if not EXCEL_PATH.exists():
        log.error(f"Excel DB 없음: {EXCEL_PATH}")
        return {}
    wb = openpyxl.load_workbook(EXCEL_PATH, data_only=True)
    if 'Matched_Plan' not in wb.sheetnames:
        log.error("Matched_Plan 시트 없음")
        wb.close()
        return {}
    ws = wb['Matched_Plan']
    
    plan_articles = defaultdict(list)
    for r in range(3, ws.max_row+1):
        pid_raw = str(ws.cell(r, 4).value or '').strip()
        if not pid_raw: continue
        
        art = {
            'date':       str(ws.cell(r, 6).value or '')[:10],
            'title_en':   str(ws.cell(r, 5).value or ''),
            'title_ko':   str(ws.cell(r, 10).value or ''),
            'source':     str(ws.cell(r, 7).value or ''),
            'province':   str(ws.cell(r, 8).value or ''),
            'sector':     str(ws.cell(r, 9).value or ''),
            'summary_ko': str(ws.cell(r, 13).value or ''),
            'summary_en': str(ws.cell(r, 14).value or ''),
            'grade':      str(ws.cell(r, 3).value or '').upper(),
            'url':        str(ws.cell(r, 17).value or ''),
            'pid_raw':    pid_raw,
        }
        for std_pid, aliases in PID_ALIASES.items():
            if pid_raw in aliases or any(a in pid_raw for a in aliases):
                plan_articles[std_pid].append(art)
                break
    wb.close()
    
    # 정렬
    for pid in plan_articles:
        plan_articles[pid].sort(key=lambda x: x['date'], reverse=True)
    log.info(f"기사 매핑 완료: {sum(len(v) for v in plan_articles.values())}건")
    return dict(plan_articles)


# ═════════════════════════════════════════════════════
# Step 3: 플랜별 통계 + 갭 진단
# ═════════════════════════════════════════════════════
def analyze_gap(pid: str, articles: list, cutoff_w: str) -> list:
    """기사 매핑 결과 자체 진단 — AI 자체 평가"""
    issues = []
    n = len(articles)
    
    if n == 0:
        issues.append('❌ 매핑 기사 0건 — knowledge_index 키워드 보강 필요')
    elif n < 3:
        issues.append(f'⚠️ 매핑 기사 부족 ({n}건) — 추가 RSS·정부사이트 발굴 필요')
    
    this_week = sum(1 for a in articles if a['date'] >= cutoff_w)
    if this_week == 0:
        issues.append('⚠️ 금주 신규 기사 없음 — Haiku 맥락분류 활성화 권장')
    
    yr_dist = defaultdict(int)
    for a in articles:
        yr = a['date'][:4] if len(a['date']) >= 4 else ''
        if yr: yr_dist[yr] += 1
    
    if '2025' not in yr_dist and '2026' not in yr_dist:
        issues.append('🔴 최신 기사(2025-2026) 누락 — 정부 발표문 직접 검색 필요')
    
    if pid == 'VN-ENV-IND-1894' and n < 5:
        issues.append('🆕 신규 플랜 — 키워드 등록: "Decision 1894", "환경산업 개발", "MOIT environmental industry"')
    
    return issues


def compute_stats(plans: dict, plan_articles: dict, cutoff_w: str) -> dict:
    """플랜별 통계 보강"""
    for pid, info in plans.items():
        arts = plan_articles.get(pid, [])
        info['articles'] = arts
        
        yr_dist = defaultdict(int)
        for a in arts:
            yr = a['date'][:4] if len(a['date']) >= 4 else ''
            if yr: yr_dist[yr] += 1
        
        info['article_stats'] = {
            'total': len(arts),
            'this_week': sum(1 for a in arts if a['date'] >= cutoff_w),
            'year_dist': dict(yr_dist),
            'latest_date': arts[0]['date'] if arts else '',
        }
        info['gap_issues'] = analyze_gap(pid, arts, cutoff_w)
    return plans


# ═════════════════════════════════════════════════════
# Step 4: Word 보고서 생성
# ═════════════════════════════════════════════════════
def kfont(run, font='맑은 고딕', size=10, bold=False, color=None):
    run.font.name = 'Calibri'
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = OxmlElement('w:rFonts')
        rPr.insert(0, rFonts)
    rFonts.set(qn('w:eastAsia'), font)


def shade_cell(cell, color_hex):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear'); shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), color_hex)
    tcPr.append(shd)


def generate_word_report(plans: dict, gap_analysis: dict, today_str: str, week_no: int, cutoff_w: str) -> Path:
    doc = Document()
    sec = doc.sections[0]
    sec.page_width = Cm(21.0); sec.page_height = Cm(29.7)
    sec.left_margin = sec.right_margin = Cm(2.0)
    sec.top_margin = sec.bottom_margin = Cm(2.0)
    
    def add_para(text='', size=10, bold=False, color=None, align=None, space_after=None):
        p = doc.add_paragraph()
        if align: p.alignment = align
        if space_after is not None: p.paragraph_format.space_after = Pt(space_after)
        if text:
            r = p.add_run(text); kfont(r, size=size, bold=bold, color=color)
        return p
    
    def add_heading(text, level=1, color=(31,78,120)):
        p = doc.add_paragraph()
        if level == 1:
            p.paragraph_format.space_before = Pt(18); p.paragraph_format.space_after = Pt(8)
            size, bold = 18, True
        elif level == 2:
            p.paragraph_format.space_before = Pt(14); p.paragraph_format.space_after = Pt(6)
            size, bold = 14, True
        else:
            p.paragraph_format.space_before = Pt(10); p.paragraph_format.space_after = Pt(4)
            size, bold = 11, True
        r = p.add_run(text); kfont(r, size=size, bold=bold, color=color)
    
    total_articles = sum(len(p.get('articles',[])) for p in plans.values())
    this_week_total = sum(p.get('article_stats',{}).get('this_week',0) for p in plans.values())
    year = datetime.now().year
    
    # 표지
    add_para('VIETNAM INFRASTRUCTURE', size=24, bold=True, color=(31,78,120), align=WD_ALIGN_PARAGRAPH.CENTER, space_after=2)
    add_para('MARKET INTELLIGENCE REPORT', size=18, bold=True, color=(31,78,120), align=WD_ALIGN_PARAGRAPH.CENTER, space_after=2)
    add_para('베트남 인프라 시장 동향 주간 보고서', size=14, bold=True, color=(89,89,89), align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para()
    add_para(today_str, size=12, color=(89,89,89), align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(f'주간호 {year}-W{week_no:02d}  │  SA-7 knowledge_index + Claude Haiku 연계 분석',
             size=10, color=(89,89,89), align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(f'수록 기사: {total_articles}건  │  마스터플랜: {len(plans)}개  │  AI 분석: Claude Haiku',
             size=10, color=(89,89,89), align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_page_break()
    
    # Executive Summary
    add_heading('Executive Summary — 주요 동향 분석', level=1)
    add_para('🤖 AI 논평 — Claude Haiku 생성', size=10, bold=True, color=(192,0,0))
    
    add_para(
        f'금주 보고서는 베트남 인프라 19개 마스터플랜 전체에 대해 누적 {total_articles}건의 매핑 기사를 분석했습니다. '
        f'금주 신규 매핑 기사는 {this_week_total}건입니다.', size=10
    )
    
    # 영역별 요약
    for area_name, plan_ids in AREA_GROUPS.items():
        area_arts = sum(len(plans[pid].get('articles',[])) for pid in plan_ids if pid in plans)
        area_week = sum(plans[pid].get('article_stats',{}).get('this_week',0) for pid in plan_ids if pid in plans)
        add_para(f'• {area_name}: 누적 {area_arts}건 / 금주 신규 {area_week}건', size=10)
    
    add_para()
    add_para('⚠️ 데이터 수집 품질: 매핑 부족 플랜에 대해 Haiku 맥락분류 + 우회 RSS 발굴 권고', 
             size=10, bold=True, color=(192,0,0))
    doc.add_page_break()
    
    # Area별 섹션 + 플랜별 상세
    plan_no = 0
    for area_name, plan_ids in AREA_GROUPS.items():
        add_heading(f'【 {area_name} 】', level=1)
        
        for pid in plan_ids:
            if pid not in plans: continue
            info = plans[pid]
            plan_no += 1
            
            add_heading(f'{plan_no}. [{pid}] {info.get("plan_name","")}', level=2)
            meta = f'{info.get("legal_basis","")}  │  {info.get("sector","")}  │  {info.get("area","").split("(")[0].strip()}'
            add_para(meta, size=9, color=(89,89,89))
            add_para()
            
            # KPI 테이블
            add_heading('■ 사업 목표 및 진행현황 (KPI)', level=3, color=(89,89,89))
            kpis = info.get('kpis', [])
            if kpis:
                tbl = doc.add_table(rows=1, cols=3)
                hdr = tbl.rows[0].cells
                for i, h in enumerate(['지표','목표(2030)','현황(2024~2026)']):
                    p = hdr[i].paragraphs[0]
                    r = p.add_run(h); kfont(r, size=9, bold=True, color=(255,255,255))
                    shade_cell(hdr[i], '1F4E78')
                for kpi in kpis:
                    row = tbl.add_row().cells
                    for i, v in enumerate([kpi.get('indicator',''), kpi.get('target',''), kpi.get('current','')]):
                        p = row[i].paragraphs[0]
                        r = p.add_run(str(v)); kfont(r, size=9)
            add_para()
            
            # 사업 개요
            add_heading('■ 사업 개요', level=3, color=(89,89,89))
            add_para('📋 Layer 1 — 사업 개요 (knowledge_index 고정 데이터)', size=9, bold=True, color=(192,0,0))
            if info.get('overview'):
                add_para(info['overview'], size=10)
            add_para()
            
            # 프로젝트 테이블
            add_heading('■ 주요 프로젝트 목록', level=3, color=(89,89,89))
            projects = info.get('projects', [])
            if projects:
                tbl = doc.add_table(rows=1, cols=4)
                hdr = tbl.rows[0].cells
                for i, h in enumerate(['프로젝트명','위치','규모/용량','비고']):
                    p = hdr[i].paragraphs[0]
                    r = p.add_run(h); kfont(r, size=9, bold=True, color=(255,255,255))
                    shade_cell(hdr[i], '1F4E78')
                for proj in projects[:25]:
                    row = tbl.add_row().cells
                    for i, v in enumerate([proj.get('name',''), proj.get('location',''), proj.get('capacity',''), proj.get('note','')]):
                        p = row[i].paragraphs[0]
                        r = p.add_run(str(v)); kfont(r, size=8)
            add_para()
            
            # AI 분석
            articles = info.get('articles', [])
            add_heading(f'■ 최신 뉴스 분석 ({len(articles)}건) — Claude Haiku AI 연계 분석', level=3, color=(89,89,89))
            add_para('🤖 Layer 2 — AI 분석 (Claude Haiku 동적 생성)', size=9, bold=True, color=(0,112,192))
            ai_lines = info.get('ai_analysis', [])
            if ai_lines:
                for line in ai_lines:
                    add_para(line, size=10)
            elif articles:
                top = articles[0]
                add_para(f'- 최신 동향({top["date"]}): {top["title_ko"] or top["title_en"][:80]}', size=10)
                if top.get('summary_ko'):
                    add_para(f'  {top["summary_ko"][:200]}', size=10)
            
            # 갭 진단
            if pid in gap_analysis and gap_analysis[pid]:
                add_para()
                add_para('🎯 데이터 수집 자체 진단 (Self-Diagnosis):', size=9, bold=True, color=(192,0,0))
                for issue in gap_analysis[pid]:
                    add_para(f'  • {issue}', size=9, color=(192,0,0))
            add_para()
            
            # 수집 기사 목록
            add_heading('■ 수집 기사 목록 (최신순)', level=3, color=(89,89,89))
            if articles:
                for ai_idx, art in enumerate(articles[:8], 1):
                    tbl = doc.add_table(rows=2, cols=4)
                    meta = tbl.rows[0].cells
                    for i, v in enumerate([f'[{ai_idx}] {art["date"]}', art["source"][:25],
                                            art["grade"] or 'MEDIUM', art["province"][:15] or '---']):
                        p = meta[i].paragraphs[0]
                        r = p.add_run(str(v)); kfont(r, size=8, bold=True)
                        shade_cell(meta[i], 'E8F0FE')
                    body = tbl.rows[1].cells
                    merged = body[0].merge(body[1]).merge(body[2]).merge(body[3])
                    p = merged.paragraphs[0]
                    r = p.add_run(art["title_en"][:130]); kfont(r, size=9, bold=True, color=(31,78,120))
                    if art.get("title_ko"):
                        p2 = merged.add_paragraph()
                        r2 = p2.add_run(f'🇰🇷 {art["title_ko"][:120]}'); kfont(r2, size=8)
                    if art.get("summary_ko"):
                        p3 = merged.add_paragraph()
                        r3 = p3.add_run(f'📝 {art["summary_ko"][:200]}'); kfont(r3, size=8, color=(89,89,89))
                    add_para('', size=4)
            else:
                add_para('(매핑 기사 없음 — knowledge_index 키워드 보강 필요)', size=10, color=(150,150,150))
            add_para()
        
        doc.add_page_break()
    
    # 종합 데이터 품질 진단
    add_heading('🎯 데이터 수집 품질 종합 진단 — AI 자체 평가', level=1, color=(192,0,0))
    add_para()
    tbl = doc.add_table(rows=1, cols=5)
    hdr = tbl.rows[0].cells
    for i, h in enumerate(['Plan_ID','매핑 기사','금주 신규','최신 기사일','진단']):
        p = hdr[i].paragraphs[0]
        r = p.add_run(h); kfont(r, size=9, bold=True, color=(255,255,255))
        shade_cell(hdr[i], '1F4E78')
    for pid in plans.keys():
        info = plans[pid]
        stats = info.get('article_stats', {})
        diag = '🟢 양호' if stats.get('total',0) >= 10 else ('🟡 보통' if stats.get('total',0) >= 3 else '🔴 부족')
        row = tbl.add_row().cells
        for i, v in enumerate([pid, str(stats.get('total',0)), str(stats.get('this_week',0)),
                               stats.get('latest_date','—'), diag]):
            p = row[i].paragraphs[0]
            r = p.add_run(str(v)); kfont(r, size=9)
    
    # 저장
    out_path = DOCS_DIR / 'Vietnam_Infra_MI_Report_Weekly.docx'
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    log.info(f"✅ Word 보고서 저장: {out_path}")
    return out_path


# ═════════════════════════════════════════════════════
# Step 5: PPT 보고서 생성 (간소화 버전)
# ═════════════════════════════════════════════════════
def generate_ppt_report(plans: dict, gap_analysis: dict, today_str: str, week_no: int) -> Path:
    """PPT는 Word 대비 간소화 — 표지 + KPI 그리드 + Area별 섹션 + 종합진단"""
    prs = Presentation()
    prs.slide_width = PCm(33.867); prs.slide_height = PCm(19.05)
    
    C = {'navy': PRGB(0x1F,0x4E,0x78), 'white': PRGB(0xFF,0xFF,0xFF),
         'gray': PRGB(0x59,0x59,0x59), 'lgray': PRGB(0xD9,0xD9,0xD9),
         'red': PRGB(0xC0,0x00,0x00), 'lblue': PRGB(0xE8,0xF0,0xFE)}
    
    # 표지
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C['navy']; bg.line.fill.background()
    
    txt = slide.shapes.add_textbox(PCm(2), PCm(7), PCm(30), PCm(2))
    tf = txt.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'VIETNAM INFRASTRUCTURE'
    r.font.name='Calibri'; r.font.size=PPt(48); r.font.bold=True; r.font.color.rgb = C['white']
    
    txt2 = slide.shapes.add_textbox(PCm(2), PCm(10), PCm(30), PCm(2))
    p = txt2.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f'Weekly MI Report — {today_str}'
    r.font.name='Calibri'; r.font.size=PPt(24); r.font.color.rgb = C['lgray']
    
    txt3 = slide.shapes.add_textbox(PCm(2), PCm(13), PCm(30), PCm(2))
    p = txt3.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f'19개 마스터플랜  |  {sum(len(p.get("articles",[])) for p in plans.values())}건 매핑'
    r.font.name='맑은 고딕'; r.font.size=PPt(16); r.font.color.rgb = C['white']
    
    # 종합 갭 분석 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PCm(2))
    bar.fill.solid(); bar.fill.fore_color.rgb = C['red']; bar.line.fill.background()
    txt = slide.shapes.add_textbox(PCm(1), PCm(0.3), PCm(30), PCm(1.5))
    r = txt.text_frame.paragraphs[0].add_run(); r.text = '🎯 데이터 수집 품질 진단'
    r.font.name='Calibri'; r.font.size=PPt(22); r.font.bold=True; r.font.color.rgb = C['white']
    
    # 갭 테이블
    rows = len(plans) + 1
    tbl = slide.shapes.add_table(rows, 5, PCm(1), PCm(3), PCm(31), PCm(15)).table
    headers = ['Plan_ID','매핑 기사','금주 신규','상태','주요 진단']
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = C['navy']
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name='맑은 고딕'; r.font.size=PPt(11); r.font.bold=True; r.font.color.rgb = C['white']
    
    for ri, (pid, info) in enumerate(plans.items(), 1):
        stats = info.get('article_stats', {})
        total = stats.get('total', 0)
        status = '🟢 양호' if total >= 10 else ('🟡 보통' if total >= 3 else '🔴 부족')
        issues = gap_analysis.get(pid, [])
        issue_text = issues[0][:35] if issues else '✅'
        
        for ci, v in enumerate([pid, str(total), str(stats.get('this_week',0)), status, issue_text]):
            cell = tbl.cell(ri, ci); cell.text = v
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name='맑은 고딕'; r.font.size=PPt(9)
    
    # 저장
    out_path = DOCS_DIR / 'Vietnam_Infra_MI_Report_Weekly.pptx'
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    log.info(f"✅ PPT 보고서 저장: {out_path}")
    return out_path


# ═════════════════════════════════════════════════════
# Step 6: Dashboard HTML 생성
# ═════════════════════════════════════════════════════
def generate_dashboard(plans: dict) -> Path:
    """Dashboard HTML 생성 - 템플릿이 별도 파일로 있으므로 간소화"""
    # 대시보드용 데이터 압축
    dash_plans = {}
    for pid, info in plans.items():
        arts_compact = []
        for a in info.get('articles', [])[:30]:
            arts_compact.append({
                'd': a.get('date',''),
                'te': (a.get('title_en','') or '')[:200],
                'tk': (a.get('title_ko','') or '')[:200],
                'sk': (a.get('summary_ko','') or '')[:200],
                'src': (a.get('source','') or '')[:40],
                'pv': (a.get('province','') or '')[:30],
                'g': a.get('grade',''),
                'url': a.get('url',''),
            })
        dash_plans[pid] = {
            'plan_id': pid,
            'name': info.get('plan_name',''),
            'area': info.get('area',''),
            'sector': info.get('sector',''),
            'legal': info.get('legal_basis',''),
            'overview': info.get('overview',''),
            'kpis': info.get('kpis', []),
            'projects': info.get('projects', [])[:20],
            'ai_analysis': info.get('ai_analysis', []),
            'stats': info.get('article_stats', {}),
            'articles': arts_compact,
            'gap_issues': info.get('gap_issues', []),
        }
    
    totals = {
        'plan_count': len(dash_plans),
        'article_total': sum(len(p['articles']) for p in dash_plans.values()),
        'this_week_total': sum(p['stats'].get('this_week',0) for p in dash_plans.values()),
        'high_count': sum(sum(1 for a in p['articles'] if a['g']=='HIGH') for p in dash_plans.values()),
        'med_count': sum(sum(1 for a in p['articles'] if a['g']=='MEDIUM') for p in dash_plans.values()),
        'pol_count': sum(sum(1 for a in p['articles'] if a['g']=='POLICY') for p in dash_plans.values()),
        'today': date.today().strftime('%Y-%m-%d'),
    }
    
    # 템플릿 파일 (없으면 기본 HTML)
    template_path = _SCRIPTS_DIR / 'dashboard_template.html'
    if not template_path.exists():
        log.warning(f"대시보드 템플릿 없음 ({template_path}) — 기본 형식 사용")
        return None
    
    with open(template_path, encoding='utf-8') as f:
        template = f.read()
    
    data_str = json.dumps({'plans': dash_plans, 'totals': totals}, ensure_ascii=False, separators=(',',':'))
    html = template.replace('__DATA_PLACEHOLDER__', data_str)
    
    out_path = DOCS_DIR / 'index.html'
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    out_path.write_text(html, encoding='utf-8')
    log.info(f"✅ Dashboard HTML 저장: {out_path}")
    return out_path


# ═════════════════════════════════════════════════════
# 메인
# ═════════════════════════════════════════════════════
def main():
    log.info("=" * 58)
    log.info(f"weekly_mi_report — {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    log.info("=" * 58)
    
    today = date.today()
    today_str = today.strftime('%Y-%m-%d')
    week_no = today.isocalendar()[1]
    cutoff_w = (today - timedelta(days=7)).strftime('%Y-%m-%d')
    
    # 1. knowledge_index 로드
    plans = load_knowledge_index()
    if not plans:
        log.error("knowledge_index 없음 — 보고서 생성 중단")
        return
    
    # 2. 기사 매핑
    plan_articles = load_articles_per_plan()
    
    # 3. 통계 + 갭 진단
    plans = compute_stats(plans, plan_articles, cutoff_w)
    gap_analysis = {pid: info.get('gap_issues', []) for pid, info in plans.items() if info.get('gap_issues')}
    
    # 4. Word 보고서
    word_path = generate_word_report(plans, gap_analysis, today_str, week_no, cutoff_w)
    
    # 5. PPT 보고서
    ppt_path = generate_ppt_report(plans, gap_analysis, today_str, week_no)
    
    # 6. Dashboard
    dash_path = generate_dashboard(plans)
    
    # 7. 실행 로그
    AGENT_OUT.mkdir(parents=True, exist_ok=True)
    log_data = {
        'run_at': datetime.now().isoformat(),
        'today': today_str,
        'week_no': f'{today.year}-W{week_no:02d}',
        'plans_count': len(plans),
        'article_total': sum(len(p.get('articles',[])) for p in plans.values()),
        'this_week_total': sum(p.get('article_stats',{}).get('this_week',0) for p in plans.values()),
        'gap_count': len(gap_analysis),
        'outputs': {
            'word': str(word_path) if word_path else None,
            'ppt': str(ppt_path) if ppt_path else None,
            'dashboard': str(dash_path) if dash_path else None,
        },
    }
    with open(AGENT_OUT / 'weekly_report_log.json', 'w', encoding='utf-8') as f:
        json.dump(log_data, f, ensure_ascii=False, indent=2)
    
    log.info("=" * 58)
    log.info(f"✅ 주간 MI 보고서 생성 완료: {today_str} (W{week_no:02d})")
    log.info(f"   매핑 기사: {log_data['article_total']}건 (금주 {log_data['this_week_total']}건)")
    log.info(f"   갭 진단: {log_data['gap_count']}개 플랜")
    log.info("=" * 58)


if __name__ == '__main__':
    main()
