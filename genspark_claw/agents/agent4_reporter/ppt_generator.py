"""
mi_ppt_generator.py — MI 보고서 PPT Executive Summary 생성기 v1.0
===================================================================
Word 보고서를 읽어 경영진 보고용 PPT 생성
  - 슬라이드 1: 표지 (플랜명, 날짜, CONFIDENTIAL)
  - 슬라이드 2: 마스터플랜 핵심 지표 (KPI 대시보드)
  - 슬라이드 3: 진행현황 요약 (프로젝트/트래킹별)
  - 슬라이드 4: 최신 뉴스 하이라이트 (상위 5건)
  - 슬라이드 5: 지역(Province)별 활동 현황
  - 슬라이드 6: 한국 기업 기회 분석
  - 슬라이드 7: 결론 및 다음 단계
"""
import os, json, ast, re
from pathlib import Path
from datetime import datetime
from collections import defaultdict, Counter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path("/home/work/claw")

# ─── 색상 팔레트 ─────────────────────────────────────────
C_NAVY   = RGBColor(0x0F, 0x17, 0x2A)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GOLD   = RGBColor(0xF5, 0x9E, 0x0B)
C_BLUE   = RGBColor(0x1E, 0x40, 0xAF)
C_LIGHT  = RGBColor(0xEF, 0xF6, 0xFF)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_GREEN  = RGBColor(0x15, 0x80, 0x3D)
C_RED    = RGBColor(0xDC, 0x26, 0x26)
C_ORANGE = RGBColor(0xEA, 0x58, 0x0C)

# ─── 플랜별 PPT 메타데이터 ────────────────────────────────
PPT_META = {
    "PDP8-INTEGRATED": {
        "title": "베트남 전력개발계획 VIII (PDP8)",
        "subtitle": "에너지 통합 보고서 — Decision 768/QD-TTg",
        "theme_rgb": RGBColor(0xD9, 0x77, 0x06),
        "accent_rgb": RGBColor(0xFF, 0xFB, 0xEB),
        "icon": "⚡",
        "kpis": [
            ("총 설비용량 2030", "183,291–236,363 MW"),
            ("태양광 목표", "46,459–73,416 MW"),
            ("원자력 재개", "4,000–6,400 MW"),
            ("투자 규모 (2026-30)", "$136.3B"),
        ],
        "highlights": [
            "Decision 768 (2025.04.15) — 원본 대비 RE 비중 대폭 상향",
            "원자력 복귀: 닌투언 1·2호기 2030-35 상업운전 목표",
            "그린수소·RE 수출: 싱가포르·말레이시아 5,000-10,000MW",
            "석탄 신규 불허, 2050년 완전 폐지 로드맵",
            "LNG 발전 22,524MW — 한국 EPC 기회 최대 $3B+",
        ],
    },
    "WATER-INTEGRATED": {
        "title": "베트남 수자원·상수도 통합 계획",
        "subtitle": "3개 마스터플랜 통합 추적 보고서",
        "theme_rgb": RGBColor(0x0E, 0xA5, 0xE9),
        "accent_rgb": RGBColor(0xF0, 0xF9, 0xFF),
        "icon": "💧",
        "kpis": [
            ("도시 상수도 보급률", "95~100% (2030)"),
            ("농촌 보급률 목표", "90% (2030)"),
            ("NRW 손실 목표", "≤15~20%"),
            ("총 투자", "$15B+ (2021-30)"),
        ],
        "highlights": [
            "3개 계획: 수자원관리(MONRE) + 도시상수도(MOC) + 농촌급수(MARD)",
            "2020 기준 폐수처리율 15% → 2030년 70% 목표",
            "6,200만 농촌인구 안전한 식수 접근 보장",
            "한국 K-water NRW 저감 기술 수출 기회",
            "EDCF ODA 패키지 $500M+ 활용 가능",
        ],
    },
    "HANOI-INTEGRATED": {
        "title": "하노이 도시개발 마스터플랜 2045/2065",
        "subtitle": "5대 도시권 통합 보고서 — Decision 1668/QD-TTg",
        "theme_rgb": RGBColor(0xDC, 0x26, 0x26),
        "accent_rgb": RGBColor(0xFE, 0xF2, 0xF2),
        "icon": "🏙️",
        "kpis": [
            ("총 투자 (2026-45)", "$2.5T"),
            ("인구 목표 2030", "12M명"),
            ("BRG 스마트시티", "$4.2B (진행중)"),
            ("메트로 목표 2035", "410km+"),
        ],
        "highlights": [
            "Decision 1668 (2024.12.27) — 5대 도시권 다핵다중심 구조",
            "북부 신도시: BRG $4.2B + 코로아 1,400ha 착공 (2026.02)",
            "서부 과학기술: 호아락 1,586ha + Tien Xuan 600ha (2026 착공)",
            "링로드4 공사중 (2023-27, 112km, $8.5B)",
            "한국 기업 접근 가능 시장: $30-50B 추정 (2026-2035)",
        ],
    },
    "VN-TRAN-2055": {
        "title": "국가 교통인프라 마스터플랜 2055",
        "subtitle": "Decision 1454/QD-TTg — $243B 투자",
        "theme_rgb": RGBColor(0xEA, 0x58, 0x0C),
        "accent_rgb": RGBColor(0xFF, 0xF7, 0xED),
        "icon": "🛣️",
        "kpis": [
            ("총 투자", "~$243B (2021-2055)"),
            ("고속도로 2030 목표", "5,000km"),
            ("Long Thanh 공항", "$16B (1단계)"),
            ("고속철도", "기본 선로 2030 목표"),
        ],
        "highlights": [
            "남북 고속도로 전 구간 완공 목표 2030",
            "Long Thanh 국제공항 1단계 — 2026 준공 목표 (지연)",
            "Lach Huyen 심수항 2단계 $2B — 아시아 3대 항만 목표",
            "하노이 링로드4 + 호치민 링로드3 동시 진행",
            "한국 EPC + ODA 패키지 기회 $5B+",
        ],
    },
    "VN-WW-2030": {
        "title": "국가 폐수처리 마스터플랜 2030",
        "subtitle": "Decision 1393/QD-TTg — 처리율 70% 목표",
        "theme_rgb": RGBColor(0x0E, 0xA5, 0xE9),
        "accent_rgb": RGBColor(0xF0, 0xF9, 0xFF),
        "icon": "🏭",
        "kpis": [
            ("처리율 현재", "~15% (2020)"),
            ("처리율 목표 2025", "20~30%"),
            ("처리율 목표 2030", "70%"),
            ("목표 처리용량", "300-400만 ㎥/일"),
        ],
        "highlights": [
            "하노이 7개 처리장 통합 120만㎥/일 — Yen So 2단계 착공",
            "호치민 9개 처리장 210만㎥/일 — Thu Duc 3단계 완공임박",
            "MBR 막기술 국산화 목표 70-80% (D1894 연계)",
            "ODA 비중 40% — EDCF·ADB·JICA 복합 패키지",
            "한국 환경건설사 EPC 기회 $2B+",
        ],
    },
    "VN-ENV-IND-1894": {
        "title": "환경산업 개발프로그램 (Decision 1894)",
        "subtitle": "MOIT 주관 — 기술 국산화 2030",
        "theme_rgb": RGBColor(0x15, 0x80, 0x3D),
        "accent_rgb": RGBColor(0xF0, 0xFD, 0xF4),
        "icon": "🌿",
        "kpis": [
            ("폐수처리 기술 국산화", "70~80% (2030)"),
            ("배기가스 처리", "60~70% (2030)"),
            ("WtE 기술", "국내 역량 구축"),
            ("총 기사 수", "448건 (역사DB)"),
        ],
        "highlights": [
            "MOIT 주관 기술산업 육성 — 인프라 건설 아님",
            "Decision 980 기술목록 등재 → 정부조달 우선 공급자 지위",
            "한국 기술이전 JV: MBR막·모니터링센서·집진기",
            "재활용 산업단지 파일럿 — 법적 근거 2027 목표",
            "MOIT 직접 접촉 채널 우선 (MONRE 아님)",
        ],
    },
    "VN-IP-NORTH-2030": {
        "title": "북부 산업단지 개발 계획 2030",
        "subtitle": "Decision 1107/QD-TTg — FDI 허브",
        "theme_rgb": RGBColor(0x7C, 0x3A, 0xED),
        "accent_rgb": RGBColor(0xF5, 0xF3, 0xFF),
        "icon": "🏭",
        "kpis": [
            ("핵심 클러스터", "5개 성시"),
            ("FDI 누적 목표", "$180B (2030)"),
            ("스마트 산업단지", "50% 전환 (2030)"),
            ("입주율 목표", "90%"),
        ],
        "highlights": [
            "삼성·LG·폭스콘 핵심 전자 제조 클러스터",
            "VSIP Hung Yen 3단계 착공 (2024)",
            "한국 산업단지 개발: LH공사·포스코인터내셔널",
            "스마트팩토리 전환 솔루션 수요 급증",
            "링로드4 완공(2027)으로 물류 혁신",
        ],
    },
    "VN-SWM-NATIONAL-2030": {
        "title": "전국 고형폐기물 통합관리 전략",
        "subtitle": "Decision 491/QD-TTg — WtE 중심 전환",
        "theme_rgb": RGBColor(0x15, 0x80, 0x3D),
        "accent_rgb": RGBColor(0xF0, 0xFD, 0xF4),
        "icon": "♻️",
        "kpis": [
            ("도시 수거율 목표", "90% (2025)"),
            ("WtE 처리 목표", "50% (2030)"),
            ("하노이 WtE", "4,000 TPD"),
            ("호치민 WtE", "2,000 TPD"),
        ],
        "highlights": [
            "매립 → 소각·WtE 전환이 국가 핵심 과제",
            "하노이 4,000TPD WtE 사업자 공모 진행 중",
            "D1894 연계 — 재활용 기술 국산화 목표",
            "GCF(녹색기후기금) + EDCF 복합 금융 가능",
            "한국 WtE: 코오롱글로벌·한국지역난방공사 컨소시엄",
        ],
    },
}
# 나머지 플랜 기본 메타
DEFAULT_META_TEMPLATE = {
    "VN-URB-METRO-2030": {"title": "도시철도·메트로 개발계획 2030", "theme_rgb": RGBColor(0xDC,0x26,0x26), "icon": "🚇",
        "kpis": [("하노이 목표","616km/15개 노선"),("호치민 목표","220km/8개 노선"),("현재 운영","25km"),("투자","$60B+")],
        "highlights": ["하노이 Line 2A·3 운영, Line 1·5·6 건설 예정","호치민 Line 1 완공 임박","현대로템 차량 공급 기회","EDCF Line 5 ODA 협의","역사 상업화 → 한국 리테일"]},
    "VN-MEKONG-DELTA-2030": {"title": "메콩델타 지역개발 마스터플랜 2030", "theme_rgb": RGBColor(0x0D,0x94,0x88), "icon": "🌊",
        "kpis": [("고속도로 목표","830km"),("현재 운영","180km"),("총 투자","$65B"),("13개 성시","1,700만명")],
        "highlights": ["기후변화 대응 — 해수 차단·염수 침투 방지","메콩 고속도로 5개 구간 2026 동시 착공","껀터 국제공항 확장","연약지반 공법 — 한국 건설사 특기"]},
    "VN-RED-RIVER-2030": {"title": "홍강 델타 지역 개발 마스터플랜 2030", "theme_rgb": RGBColor(0xBE,0x12,0x3C), "icon": "🏔️",
        "kpis": [("11개 성시","2,300만명"),("고속도로 목표","1,300km"),("FDI 목표 2030","$15B/년"),("총 투자","$150B+")],
        "highlights": ["하노이-하이퐁-꽝닌 삼각축 경제권","삼성·LG 전자 클러스터 허브","링로드4 착공 112km/$8.5B","한국 소부장 동반 진출 최적지"]},
    "VN-SC-2030": {"title": "스마트시티 국가전략 2030 (기술·플랫폼)", "theme_rgb": RGBColor(0x7C,0x3A,0xED), "icon": "🏙️",
        "kpis": [("플랫폼 목표","63개 성시"),("전자정부","ASEAN Top 3"),("투자","$15B"),("5G 커버","전국")],
        "highlights": ["기술 레이어 — 공간 레이어(Hanoi 보고서)와 구분","AI 교통관제·디지털트윈·IoT 인프라","SK텔레콤·KT·NIA 솔루션 수출","KOICA ODA 스마트시티 패키지"]},
    "VN-OG-2030": {"title": "석유가스 개발 계획 2030", "theme_rgb": RGBColor(0xD9,0x77,0x06), "icon": "⛽",
        "kpis": [("원유 생산","800만톤/년"),("가스 생산","130억㎥/년"),("LNG 터미널","5개소+"),("탐사 블록","15+ 블록")],
        "highlights": ["Bach Ho 고갈 → 심해 탐사 긴급","Thi Vai LNG 2단계 + Ca Mau LNG 건설","GS에너지·SK이노베이션 블록 JV 기회","KOGAS 운영 노하우 수출"]},
    "VN-EV-2030": {"title": "전기차·친환경 모빌리티 2030", "theme_rgb": RGBColor(0x0D,0x94,0x88), "icon": "🚗",
        "kpis": [("EV 판매비중 2030","30%"),("충전소 목표","10만 기"),("전기버스 전환","50% (2030)"),("배터리","국내 3곳+")],
        "highlights": ["VinFast 급성장 — 2024 EV 점유율 8%→30% 목표","삼성SDI·LG에너지솔루션 VinFast 배터리 공급","현대차 Elec City 전기버스 입찰","포스코퓨처엠 양극재 JV 추진"]},
    "VN-CARBON-2050": {"title": "탄소중립 2050 로드맵", "theme_rgb": RGBColor(0x16,0x65,0x34), "icon": "🌍",
        "kpis": [("GHG 감축 2030","43.5%"),("탄소중립","2050"),("ETS 공식 운영","2028"),("탄소시장","$500M+ (2030)"),],
        "highlights": ["COP26 공약 이행 — NDC 2022 업데이트","ETS 시범 2025 → 공식 2028","K-ETS 운영 경험 → 베트남 ETS 설계 자문","REDD+·산림 크레딧 한국 기업 확보 기회"]},
}

def _is_vi(text):
    vi_chars = set('ăâêôơưđáàảãạắằẳẵặấầẩẫậéèẻẽẹếềểễệíìỉĩịóòỏõọốồổỗộớờởỡợúùủũụứừửữựýỳỷỹỵ')
    return any(c.lower() in vi_chars for c in (text or ''))

def _date_key(art):
    d = str(art.get('published_date','') or '')
    try:
        if re.match(r'\d{4}-\d{2}-\d{2}', d): return d[:10]
        return datetime.strptime(d.strip(), '%b %d, %Y').strftime('%Y-%m-%d')
    except: return '1900-01-01'

def _best_title(art):
    t = art.get('title','') or ''
    if not _is_vi(t): return t[:100]
    en = art.get('summary_en','') or ''
    if en and not _is_vi(en): return f"[VI→EN] {en[:90]}"
    return t[:100]

def _best_summary_short(art, max_len=140):
    ko = art.get('summary_ko','') or ''
    en = art.get('summary_en','') or ''
    if ko and not _is_vi(ko): return ko[:max_len]
    if en and not _is_vi(en): return en[:max_len]
    return ''

def _parse_plans(art):
    mp = art.get('matched_plans',[]) or []
    if isinstance(mp, str):
        try: mp = ast.literal_eval(mp)
        except: mp = []
    return mp

# ─── PPT 빌더 헬퍼 ───────────────────────────────────────
def _add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def _txb(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def _set_text(txb, text, size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tf = txb.text_frame; tf.word_wrap = wrap
    tf.text = ''
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color

def _add_run(para, text, size=14, bold=False, color=C_WHITE):
    run = para.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    return run

def _rect(slide, left, top, width, height, fill_rgb, alpha=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

# ─── 슬라이드 생성 함수들 ────────────────────────────────

def slide_cover(prs, meta, plan_id, now, week):
    slide = _add_slide(prs, 6)
    theme = meta.get("theme_rgb", C_BLUE)
    icon = meta.get("icon", "📊")

    # 배경 전체 네이비
    _rect(slide, 0, 0, 10, 7.5, C_NAVY)
    # 하단 강조 바
    _rect(slide, 0, 6.5, 10, 1.0, theme)

    # 브랜드 텍스트
    t1 = _txb(slide, 0.4, 0.3, 9.2, 0.4)
    _set_text(t1, "VIETNAM INFRASTRUCTURE INTELLIGENCE HUB", 9, False, RGBColor(0x93,0xC5,0xFD), PP_ALIGN.CENTER)

    # 아이콘
    ti = _txb(slide, 4.0, 1.0, 2.0, 1.0)
    _set_text(ti, icon, 48, False, theme, PP_ALIGN.CENTER)

    # 제목
    t2 = _txb(slide, 0.5, 2.1, 9.0, 1.0)
    _set_text(t2, meta.get("title",""), 28, True, theme, PP_ALIGN.CENTER)

    # 부제
    t3 = _txb(slide, 0.5, 3.2, 9.0, 0.5)
    _set_text(t3, meta.get("subtitle",""), 14, False, RGBColor(0xCB,0xD5,0xE1), PP_ALIGN.CENTER)

    # 구분선
    _rect(slide, 1.5, 3.85, 7.0, 0.03, theme)

    # 날짜/주차
    t4 = _txb(slide, 0.5, 4.0, 9.0, 0.4)
    _set_text(t4, f"W{week:02d} / {now.strftime('%Y년 %m월 %d일')}  ·  CONFIDENTIAL", 11, False,
              RGBColor(0x94,0xA3,0xB8), PP_ALIGN.CENTER)

    # 하단 배너
    t5 = _txb(slide, 0.3, 6.55, 9.4, 0.4)
    _set_text(t5, f"Plan ID: {plan_id}  ·  생성: {now.strftime('%Y-%m-%d %H:%M')} UTC  ·  경영진 보고용", 9,
              False, C_NAVY, PP_ALIGN.CENTER)

def slide_kpi_dashboard(prs, meta, art_count):
    slide = _add_slide(prs, 6)
    theme = meta.get("theme_rgb", C_BLUE)
    kpis = meta.get("kpis", [])

    _rect(slide, 0, 0, 10, 1.1, C_NAVY)
    t = _txb(slide, 0.3, 0.2, 9.4, 0.7)
    _set_text(t, f"📊  핵심 KPI 대시보드  —  {meta.get('title','')}", 18, True, C_WHITE)

    # KPI 박스 4개
    cols = min(4, len(kpis))
    box_w = 9.2 / cols
    for i, (label, value) in enumerate(kpis[:4]):
        x = 0.4 + i * box_w
        # 박스 배경
        _rect(slide, x, 1.3, box_w - 0.1, 1.6, theme)
        # 레이블
        tl = _txb(slide, x + 0.05, 1.35, box_w - 0.2, 0.35)
        _set_text(tl, label, 10, False, RGBColor(0xFF,0xFF,0xCC), PP_ALIGN.CENTER)
        # 값
        tv = _txb(slide, x + 0.05, 1.7, box_w - 0.2, 0.9)
        _set_text(tv, value, 14, True, C_WHITE, PP_ALIGN.CENTER)

    # 기사 건수 박스
    _rect(slide, 0.4, 3.1, 2.0, 0.9, RGBColor(0x1E, 0x40, 0xAF))
    tc = _txb(slide, 0.5, 3.15, 1.8, 0.8)
    _set_text(tc, f"관련 기사\n{art_count}건", 13, True, C_WHITE, PP_ALIGN.CENTER)

    # 하이라이트 목록
    highlights = meta.get("highlights", [])
    th = _txb(slide, 2.6, 3.1, 7.1, 3.8)
    tf = th.text_frame; tf.word_wrap = True
    for i, hl in enumerate(highlights[:5]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▶  {hl}"
        run.font.size = Pt(12)
        run.font.color.rgb = C_NAVY if i % 2 == 0 else RGBColor(0x1E, 0x40, 0xAF)

def slide_project_status(prs, meta, arts):
    slide = _add_slide(prs, 6)
    theme = meta.get("theme_rgb", C_BLUE)

    _rect(slide, 0, 0, 10, 1.1, theme)
    t = _txb(slide, 0.3, 0.2, 9.4, 0.7)
    _set_text(t, "🚧  주요 프로젝트 진행현황 추적", 18, True, C_WHITE)

    # 2026 최신 기사 상위 5건 요약
    recent = sorted([a for a in arts if "2026" in str(a.get('published_date',''))],
                    key=_date_key, reverse=True)[:5]
    if not recent:
        recent = sorted(arts, key=_date_key, reverse=True)[:5]

    y_start = 1.3
    for i, art in enumerate(recent):
        title = _best_title(art)[:75]
        sm = _best_summary_short(art, 110)
        date = _date_key(art)
        src = (art.get('source','') or '')[:20]

        # 번호 박스
        _rect(slide, 0.3, y_start + i*1.1, 0.35, 0.9, theme)
        tn = _txb(slide, 0.31, y_start + i*1.1 + 0.2, 0.33, 0.5)
        _set_text(tn, str(i+1), 14, True, C_WHITE, PP_ALIGN.CENTER)

        # 제목
        tt = _txb(slide, 0.75, y_start + i*1.1, 8.9, 0.4)
        _set_text(tt, title, 11, True, C_NAVY)

        # 요약 + 날짜
        ts = _txb(slide, 0.75, y_start + i*1.1 + 0.42, 8.9, 0.35)
        _set_text(ts, f"{'  ' + sm[:100] if sm else ''}   📅 {date}  |  {src}", 9, False, C_GRAY)

def slide_news_highlights(prs, meta, arts):
    """최신 뉴스 하이라이트 (경영진 요약 카드형)"""
    slide = _add_slide(prs, 6)
    theme = meta.get("theme_rgb", C_BLUE)

    _rect(slide, 0, 0, 10, 1.1, C_NAVY)
    t = _txb(slide, 0.3, 0.2, 9.4, 0.7)
    _set_text(t, "📰  최신 뉴스 하이라이트 (2026년)", 18, True, C_WHITE)

    # 연도별 기사 수 바차트 (텍스트형)
    from collections import Counter
    yr_cnt = Counter()
    for a in arts:
        yr = _date_key(a)[:4]
        yr_cnt[yr] += 1

    yt = _txb(slide, 0.3, 1.15, 9.4, 0.45)
    tf = yt.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    for yr in sorted(yr_cnt.keys())[-5:]:
        cnt = yr_cnt[yr]
        bar = "▐" * min(cnt // 5, 15)
        run = p.add_run()
        run.text = f"  {yr}: {cnt:3d}건  {bar}"
        run.font.size = Pt(9); run.font.color.rgb = C_GRAY

    # 뉴스 카드 3×2
    recent = sorted([a for a in arts if "2026" in str(a.get('published_date',''))],
                    key=_date_key, reverse=True)[:6]
    if not recent:
        recent = sorted(arts, key=_date_key, reverse=True)[:6]

    for i, art in enumerate(recent):
        col = i % 3; row = i // 3
        x = 0.25 + col * 3.22; y = 1.75 + row * 2.55
        # 카드 배경
        bg = theme if i == 0 else RGBColor(0xF1, 0xF5, 0xF9)
        _rect(slide, x, y, 3.0, 2.3, bg)
        # 제목
        title = _best_title(art)[:65]
        tt = _txb(slide, x + 0.08, y + 0.05, 2.85, 0.9)
        color = C_WHITE if i == 0 else C_NAVY
        _set_text(tt, title, 10, True, color)
        # 요약
        sm = _best_summary_short(art, 100)
        ts = _txb(slide, x + 0.08, y + 0.95, 2.85, 1.05)
        _set_text(ts, sm, 9, False, C_WHITE if i == 0 else C_GRAY)
        # 날짜
        td = _txb(slide, x + 0.08, y + 2.0, 2.85, 0.25)
        _set_text(td, f"📅 {_date_key(art)}", 8, False,
                  RGBColor(0xFF,0xFF,0xCC) if i==0 else C_GRAY)

def slide_province_activity(prs, meta, arts):
    """Province 활동 현황"""
    slide = _add_slide(prs, 6)
    theme = meta.get("theme_rgb", C_BLUE)

    _rect(slide, 0, 0, 10, 1.1, theme)
    t = _txb(slide, 0.3, 0.2, 9.4, 0.7)
    _set_text(t, "🗺️  지역별 활동 현황 분석", 18, True, C_WHITE)

    ALIASES = {
        "vietnam":"🇻🇳 전국","national":"🇻🇳 전국","ho chi minh":"Ho Chi Minh City",
        "hcm":"Ho Chi Minh City","hanoi":"Hanoi","hà nội":"Hanoi","da nang":"Da Nang",
        "hai phong":"Hai Phong","can tho":"Can Tho","binh duong":"Binh Duong",
        "dong nai":"Dong Nai","quang ninh":"Quang Ninh",
    }
    counter = Counter()
    for art in arts:
        prov = (art.get('province') or art.get('location') or 'Unknown').lower().strip()
        for alias, norm in ALIASES.items():
            if alias in prov: prov = norm; break
        counter[prov.title() if len(prov) < 20 else prov[:20]] += 1

    top_provs = counter.most_common(10)
    max_cnt = top_provs[0][1] if top_provs else 1

    # 왼쪽: 바차트
    for i, (prov, cnt) in enumerate(top_provs[:8]):
        y = 1.3 + i * 0.68
        bar_w = (cnt / max_cnt) * 4.5
        _rect(slide, 0.3, y, bar_w, 0.45, theme)
        tl = _txb(slide, 0.32, y + 0.08, bar_w - 0.1, 0.3)
        _set_text(tl, f"{prov[:18]} ({cnt}건)", 9, True, C_WHITE if bar_w > 0.5 else C_NAVY)

    # 오른쪽: 연도별 트렌드
    from collections import defaultdict
    yr_prov = defaultdict(Counter)
    for art in arts:
        yr = _date_key(art)[:4]
        prov = (art.get('province') or 'Unknown')[:15]
        yr_prov[yr][prov] += 1

    t2 = _txb(slide, 5.0, 1.3, 4.7, 0.4)
    _set_text(t2, "연도별 기사 트렌드", 13, True, C_NAVY)
    tf2 = _txb(slide, 5.0, 1.75, 4.7, 5.2)
    tfr = tf2.text_frame; tfr.word_wrap = False
    for yr in sorted(yr_prov.keys())[-5:]:
        cnt = sum(yr_prov[yr].values())
        bar = "█" * min(cnt // 10, 18)
        p2 = tfr.add_paragraph()
        p2.space_before = Pt(3)
        run = p2.add_run()
        run.text = f"{yr}: {cnt:3d}건  {bar}"
        run.font.size = Pt(11); run.font.color.rgb = C_NAVY

def slide_korean_opportunity(prs, meta):
    """한국 기업 기회 분석"""
    slide = _add_slide(prs, 6)
    theme = meta.get("theme_rgb", C_BLUE)

    _rect(slide, 0, 0, 10, 1.1, RGBColor(0x0F, 0x17, 0x2A))
    t = _txb(slide, 0.3, 0.2, 9.4, 0.7)
    _set_text(t, "🇰🇷  한국 기업 기회 분석", 18, True, C_WHITE)

    # 기회 목록 (meta에서 로드, 없으면 기본값)
    from report_kpi import KPI_CONFIG
    from report_regional import PLAN_CONFIG

    plan_id_raw = meta.get("_plan_id","")
    opps = []
    if plan_id_raw in KPI_CONFIG:
        opps = KPI_CONFIG[plan_id_raw].get("korean_opportunities",[])
    elif plan_id_raw in PLAN_CONFIG:
        opps = PLAN_CONFIG[plan_id_raw].get("korean_opportunities",[])

    if not opps:
        opps = [
            ("전략적 진입 시점","HIGH","2026년 대규모 착공 시작 — 지금이 최적"),
            ("ODA 패키지","HIGH","EDCF(수출입은행) + KOICA 복합 활용"),
            ("기술이전 파트너십","MEDIUM","한국 기술 → 현지 생산 JV 구조"),
        ]

    colors = {"HIGH": RGBColor(0xDC,0x26,0x26), "MEDIUM": RGBColor(0xEA,0x58,0x0C),
              "LOW": RGBColor(0xCA,0x8A,0x04)}

    for i, (cat, level, desc) in enumerate(opps[:5]):
        y = 1.25 + i * 1.18
        bg = colors.get(level, C_BLUE)
        # 레벨 박스
        _rect(slide, 0.3, y, 1.1, 0.9, bg)
        tl = _txb(slide, 0.31, y + 0.2, 1.08, 0.5)
        _set_text(tl, level, 11, True, C_WHITE, PP_ALIGN.CENTER)
        # 카테고리
        _rect(slide, 1.5, y, 1.8, 0.9, RGBColor(0xF1, 0xF5, 0xF9))
        tc = _txb(slide, 1.55, y + 0.15, 1.7, 0.6)
        _set_text(tc, cat, 11, True, C_NAVY)
        # 설명
        td = _txb(slide, 3.4, y + 0.08, 6.3, 0.75)
        _set_text(td, desc[:110], 10, False, C_NAVY)

def slide_conclusion(prs, meta, plan_id, total_arts, week):
    """결론 및 다음 단계"""
    slide = _add_slide(prs, 6)
    theme = meta.get("theme_rgb", C_BLUE)

    _rect(slide, 0, 0, 10, 1.1, theme)
    t = _txb(slide, 0.3, 0.2, 9.4, 0.7)
    _set_text(t, "✅  결론 및 다음 단계", 18, True, C_WHITE)

    # 핵심 takeaway
    tt = _txb(slide, 0.4, 1.3, 9.2, 0.5)
    _set_text(tt, f"📌  이번 주 핵심 메시지 — {meta.get('title','')}", 14, True, C_NAVY)

    highlights = meta.get("highlights", [])
    th = _txb(slide, 0.4, 1.9, 9.2, 2.2)
    tf = th.text_frame; tf.word_wrap = True
    for i, hl in enumerate(highlights[:3]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"▶  {hl}"
        run.font.size = Pt(12); run.font.color.rgb = C_NAVY

    # 다음 단계
    _rect(slide, 0.4, 4.2, 9.2, 0.4, RGBColor(0xEF, 0xF6, 0xFF))
    tn = _txb(slide, 0.5, 4.25, 9.0, 0.3)
    _set_text(tn, "📋  다음 단계  (Next Actions)", 12, True, C_BLUE)

    actions = [
        "다음 주 수집 기사와 진행현황 비교 — 주간 트래킹 지속",
        f"관련 기사 총 {total_arts}건 — Excel DB 및 Word 보고서 상세 참조",
        "한국 기업 접촉: KOTRA 하노이 + 대한상의 베트남 활용",
    ]
    ta = _txb(slide, 0.4, 4.65, 9.2, 1.6)
    taf = ta.text_frame; taf.word_wrap = True
    for i, act in enumerate(actions):
        p = taf.add_paragraph() if i > 0 else taf.paragraphs[0]
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  {i+1}.  {act}"
        run.font.size = Pt(11); run.font.color.rgb = C_NAVY

    # 하단
    _rect(slide, 0, 7.1, 10, 0.4, C_NAVY)
    tf2 = _txb(slide, 0.3, 7.14, 9.4, 0.28)
    _set_text(tf2, f"Plan ID: {plan_id}  |  W{week:02d}  |  Vietnam Infrastructure Intelligence Hub  |  CONFIDENTIAL",
              8, False, RGBColor(0x94, 0xA3, 0xB8), PP_ALIGN.CENTER)


# ─── 메인 생성 함수 ──────────────────────────────────────
def generate_ppt(plan_id: str, output_dir: str) -> str:
    """단일 플랜 PPT 생성"""
    import sys; sys.path.insert(0, str(BASE_DIR/"scripts"))

    now = datetime.now(); week = now.isocalendar()[1]

    # 메타 로드
    meta = PPT_META.get(plan_id) or DEFAULT_META_TEMPLATE.get(plan_id)
    if not meta:
        # 기본 메타 생성
        meta = {"title": plan_id, "subtitle": "", "theme_rgb": C_BLUE, "icon": "📊",
                "kpis": [], "highlights": []}
    meta["_plan_id"] = plan_id

    # 기사 로드
    from report_lib import load_history, _date_key as date_key
    buckets = load_history([plan_id])
    arts = sorted(buckets.get(plan_id, []), key=_date_key, reverse=True)

    # PPT 생성
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # 7개 슬라이드
    slide_cover(prs, meta, plan_id, now, week)
    slide_kpi_dashboard(prs, meta, len(arts))
    slide_project_status(prs, meta, arts)
    slide_news_highlights(prs, meta, arts)
    slide_province_activity(prs, meta, arts)
    slide_korean_opportunity(prs, meta)
    slide_conclusion(prs, meta, plan_id, len(arts), week)

    # 저장
    safe_id = plan_id.replace("-","_")
    fname = f"MI_PPT_{safe_id}_W{week:02d}_{now.strftime('%Y%m%d')}.pptx"
    fpath = os.path.join(output_dir, fname)
    prs.save(fpath)
    return fpath


def generate_all_ppts(output_dir: str) -> list:
    """전체 15개 플랜 PPT 생성"""
    all_plan_ids = list(PPT_META.keys()) + [k for k in DEFAULT_META_TEMPLATE if k not in PPT_META]
    results = []
    for plan_id in all_plan_ids:
        print(f"  📊 {plan_id} ...", end=" ", flush=True)
        try:
            fpath = generate_ppt(plan_id, output_dir)
            sz = os.path.getsize(fpath) // 1024
            print(f"✅ {sz}KB")
            results.append((plan_id, fpath, sz))
        except Exception as e:
            print(f"❌ {e}")
            results.append((plan_id, None, 0))
    return results


if __name__ == "__main__":
    import sys
    sys.path.insert(0, str(BASE_DIR/"scripts"))
    out = str(BASE_DIR/"outputs/reports/MI_PPT")
    Path(out).mkdir(parents=True, exist_ok=True)
    plan = sys.argv[1] if len(sys.argv) > 1 else "PDP8-INTEGRATED"
    print(f"PPT 생성: {plan}")
    fpath = generate_ppt(plan, out)
    print(f"✅ {fpath}  ({os.path.getsize(fpath)//1024}KB)")
